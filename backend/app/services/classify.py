"""Inbox-Verarbeitung: Dateien klassifizieren, Text extrahieren, ablegen.
Migriert aus _agent/heartbeat.py. Als importierbares Modul statt Subprocess-Aufruf,
damit der Chat-Endpoint (/api/inbox_process) es direkt callen kann statt einen
neuen Python-Prozess zu starten.
"""
import hashlib
import json
import logging
import re
import shutil
import time
from datetime import date, datetime
from pathlib import Path

from app.config import get_settings
from app.constants import Models
from app.services.anthropic_client import complete_json

SKIP_EXTENSIONS = {
    ".js", ".ts", ".map", ".mjs", ".jsx", ".tsx", ".css",
    ".yml", ".yaml", ".eslintrc", ".nycrc", ".npmignore",
    ".enc", ".lock", ".editorconfig", ".orig",
}
SKIP_NAMES = {
    "Thumbs.db", ".DS_Store", "package-lock.json", "package.json",
    "tsconfig.json", "tsdoc-metadata.json", "openai",
}
SKIP_PREFIXES = ("README", "readme", "LICENSE", "license", "CHANGELOG",
                  "HISTORY", "History", "CONTRIBUTING", "contributing",
                  # Word/Excel-Sperrdateien ("~$Bericht.docx") - enthalten
                  # keinen Dokumentinhalt und scheiterten deshalb immer.
                  "~$")
logger = logging.getLogger("brain.classify")

IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".mp4", ".mov", ".mp3"}

# Textmarker, an denen ein Gesprächstranskript zuverlässig erkennbar ist -
# Teams exportiert mit "...-Besprechungstranskript" im Kopf und "Transkription
# gestartet/beendet" als Rahmen, Google Meet mit "New Record from Google Meet".
TRANSCRIPT_MARKERS = (
    "besprechungstranskript",
    "transkription gestartet",
    "meeting transcript",
    "new record from google meet",
)


def _cache_path() -> Path:
    return get_settings().agent_dir / "logs" / "processed_cache.json"


def _load_cache() -> set:
    path = _cache_path()
    if path.exists():
        try:
            return set(json.loads(path.read_text(encoding="utf-8")))
        except Exception:
            return set()
    return set()


def _save_cache(cache: set) -> None:
    _cache_path().write_text(json.dumps(list(cache)), encoding="utf-8")


# ── Dedup-Check (Klassifizierungs-Robustheit, 2026-08-20) ──────────────────
# Beobachtetes Problem: dieselbe Quelldatei kam über verschiedene Wege
# (erneuter E-Mail-Anhang-Download, doppelter Chat-Upload) wiederholt in
# _inbox/ an und wurde jedes Mal neu klassifiziert - u.a. 41 identische
# Dateinamen gleichzeitig in _inbox/_fehler/ UND _inbox/_verarbeitet/. Hash
# über die Rohbytes (nicht den extrahierten Text), damit auch bei
# unterschiedlichem OCR-Ergebnis dasselbe Original erkannt wird.

def _file_hash(filepath: Path) -> str:
    h = hashlib.sha256()
    with open(filepath, "rb") as f:
        for chunk in iter(lambda: f.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()


def _content_hash_path() -> Path:
    return get_settings().agent_dir / "logs" / "content_hashes.json"


def _load_content_hashes() -> dict:
    path = _content_hash_path()
    if path.exists():
        try:
            return json.loads(path.read_text(encoding="utf-8"))
        except Exception:
            return {}
    # Einmaliger Bootstrap: baut die Karte aus den bereits archivierten
    # Originalen in _inbox/_verarbeitet/ auf, damit der Dedup-Check auch die
    # schon bekannte Alt-Situation erkennt statt nur künftige Duplikate.
    hashes: dict = {}
    verarbeitet = get_settings().inbox_dir / "_verarbeitet"
    if verarbeitet.exists():
        for f in verarbeitet.iterdir():
            if f.is_file():
                try:
                    hashes[_file_hash(f)] = f"_inbox/_verarbeitet/{f.name}"
                except Exception:
                    continue
    _save_content_hashes(hashes)
    return hashes


def _save_content_hashes(hashes: dict) -> None:
    _content_hash_path().parent.mkdir(parents=True, exist_ok=True)
    _content_hash_path().write_text(json.dumps(hashes, ensure_ascii=False, indent=2), encoding="utf-8")


def _move_aside(filepath: Path, ziel_ordner: Path) -> None:
    """Verschiebt filepath nach ziel_ordner, durchnummeriert bei Kollision.
    Für Dedup-/Trivial-Bild-Fälle, die NICHT nach _fehler/ sollen."""
    ziel_ordner.mkdir(parents=True, exist_ok=True)
    ziel = ziel_ordner / filepath.name
    zaehler = 1
    while ziel.exists():
        ziel = ziel_ordner / f"{filepath.stem}({zaehler}){filepath.suffix}"
        zaehler += 1
    try:
        shutil.move(str(filepath), str(ziel))
    except Exception:
        pass


# ── Trivial-Bild-Filter (Klassifizierungs-Robustheit, 2026-08-20) ──────────
# extract_text() liefert für JEDES Bild denselben Platzhalter - auch für
# winzige E-Mail-Icons/Signaturbilder. Die liefen bisher trotzdem durch einen
# vollen LLM-Klassifizierungs-Call und landeten dauerhaft im Vault (u.a.
# Memos/Medien/Icons|EmailSignaturen|Logos, "image001.png" in Finanzen/).
# Bewusst konservativ (nur sehr kleine Dateien / typische Outlook-Inline-
# Namen) - im Zweifel lieber durch die normale Pipeline laufen lassen als ein
# echtes Foto/Scan zu verpassen. Nichts wird gelöscht, nur nach
# _agent/trivial_images/ verschoben (Vault-Regel: nie ohne Bestätigung
# löschen).
_TRIVIAL_IMAGE_MAX_BYTES = 20_000
_TRIVIAL_IMAGE_NAME_RE = re.compile(r"^image\d{3,}\.(png|jpe?g|gif)$", re.IGNORECASE)


def _is_trivial_image(filepath: Path) -> bool:
    if filepath.suffix.lower() not in {".png", ".jpg", ".jpeg", ".gif", ".webp"}:
        return False
    try:
        size = filepath.stat().st_size
    except OSError:
        return False
    if _TRIVIAL_IMAGE_NAME_RE.match(filepath.name) and size < 100_000:
        return True
    return size < _TRIVIAL_IMAGE_MAX_BYTES


def _sanitize_zielordner(raw: str) -> str:
    """Säubert/validiert den vom LLM gelieferten Zielordner-Pfad, bevor er als
    echter Dateisystempfad verwendet wird - vorher wurde result['zielordner']
    ungeprüft übernommen. Jedes Pfadsegment läuft durch _ordnername()
    (Länge/Sonderzeichen), '..'-Segmente und leere/absolute Pfade fallen auf
    Memos/ zurück - das Ergebnis bleibt dadurch garantiert innerhalb des
    Vaults."""
    raw = (raw or "").strip().replace("\\", "/")
    if raw.startswith("/"):
        return "Memos"
    parts = [p for p in raw.split("/") if p and p != "."]
    if not parts or any(p == ".." for p in parts):
        return "Memos"
    return "/".join(_ordnername(p) for p in parts)


def _extract_pdf_via_mistral_ocr(filepath: Path, max_chars: int) -> str | None:
    """OCR-Vorstufe für PDFs (Umsetzungsplan-Memo 2026-07-16, Punkt C1): Mistral
    liest gescannte/mehrspaltige PDFs zuverlässiger als die reine PyPDF2-
    Textextraktion aus - genau das Problem, das in diesem Projekt wiederholt zu
    abgeschnittenen/lückenhaften .md-Extrakten geführt hat. Gibt None zurück,
    wenn irgendetwas schiefgeht (kein Key, Netzwerkfehler, unerwartete Antwort) -
    der Aufrufer fällt dann automatisch auf PyPDF2 zurück, siehe extract_text()."""
    settings = get_settings()
    if not settings.mistral_api_key:
        return None
    try:
        import base64
        import requests

        b64 = base64.b64encode(filepath.read_bytes()).decode("ascii")
        resp = requests.post(
            "https://api.mistral.ai/v1/ocr",
            headers={"Authorization": f"Bearer {settings.mistral_api_key}"},
            json={
                "model": "mistral-ocr-latest",
                "document": {
                    "type": "document_url",
                    "document_url": f"data:application/pdf;base64,{b64}",
                },
            },
            timeout=60,
        )
        resp.raise_for_status()
        pages = resp.json().get("pages", [])
        text = "\n\n".join(p.get("markdown", "") for p in pages).strip()
        return text[:max_chars] if text else None
    except Exception:
        return None


def extract_text(filepath: Path, max_chars: int = 3000) -> str | None:
    suffix = filepath.suffix.lower()
    try:
        if suffix == ".pdf":
            ocr_text = _extract_pdf_via_mistral_ocr(filepath, max_chars)
            if ocr_text:
                return ocr_text

            import PyPDF2

            with open(filepath, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                return " ".join(page.extract_text() or "" for page in reader.pages)[:max_chars]
        if suffix == ".docx":
            from docx import Document

            doc = Document(filepath)
            # Absätze UND Tabellenzellen (korrigiert 2026-08-11): python-docx
            # liefert über doc.paragraphs ausschließlich Fließtext-Absätze,
            # alles in Tabellen bleibt unsichtbar. Transkript-Exporte (Notta &
            # Co.) legen Sprecher/Zeitstempel/Text aber typischerweise als
            # Tabelle an - für solche Dateien kam hier ein praktisch leerer
            # String heraus, die anschließende Klassifizierung bekam nichts zu
            # sehen und die Datei landete in _inbox/_fehler/.
            parts = [p.text for p in doc.paragraphs]
            for table in doc.tables:
                # row.cells gibt verbundene Zellen einmal pro überspanntem
                # Feld zurück; über die zugrundeliegende XML-Zelle
                # deduplizieren, sonst steht ihr Text mehrfach im Extrakt.
                seen_cells: set[int] = set()
                for row in table.rows:
                    for cell in row.cells:
                        marker = id(cell._tc)
                        if marker in seen_cells:
                            continue
                        seen_cells.add(marker)
                        parts.append(cell.text)
            return " ".join(t for t in parts if t and t.strip())[:max_chars]
        if suffix in (".xlsx", ".xls"):
            import openpyxl

            wb = openpyxl.load_workbook(filepath, read_only=True)
            text = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    text.append(" ".join(str(c) for c in row if c))
            return " ".join(text)[:max_chars]
        if suffix in (".pptx", ".ppt"):
            from pptx import Presentation

            prs = Presentation(filepath)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return " ".join(text)[:max_chars]
        if suffix == ".html":
            from html.parser import HTMLParser

            class TextExtractor(HTMLParser):
                def __init__(self):
                    super().__init__()
                    self.parts: list[str] = []

                def handle_data(self, data):
                    self.parts.append(data)

            p = TextExtractor()
            p.feed(filepath.read_text(encoding="utf-8", errors="ignore"))
            return " ".join(p.parts)[:max_chars]
        if suffix == ".json":
            raw = filepath.read_text(encoding="utf-8", errors="ignore")[:max_chars]
            return f"[JSON-Datei] {raw}"
        if suffix in (".txt", ".md", ".markdown", ".csv"):
            return filepath.read_text(encoding="utf-8", errors="ignore")[:max_chars]
        if suffix in (".png", ".jpg", ".jpeg", ".gif", ".svg", ".webp"):
            return "[Bilddatei, kein Text extrahierbar]"
        if suffix in (".mp4", ".mov", ".avi", ".mkv", ".mp3"):
            return "[Mediendatei, kein Text extrahierbar]"
        if suffix == ".zip":
            return "[ZIP-Archiv]"
        return f"[Unbekanntes Format: {suffix}]"
    except Exception:
        return None


def list_customer_names() -> list[str]:
    """Ordnernamen unter Kunden/ (ohne Platzhalter/Vorlage) - Basis für die
    automatische Zuordnung von E-Mails zu bestehenden Kunden."""
    kunden_dir = get_settings().vault_path / "Kunden"
    if not kunden_dir.exists():
        return []
    return [
        p.name for p in sorted(kunden_dir.iterdir())
        if p.is_dir() and not p.name.startswith((".", "[", "_"))
    ]


def list_lead_names() -> list[str]:
    """Dateinamen (ohne Datumspräfix) unter Leads/ - analoge Basis für die
    automatische Zuordnung von E-Mails zu Leads (siehe list_customer_names())."""
    leads_dir = get_settings().vault_path / "Leads"
    if not leads_dir.exists():
        return []
    return [
        re.sub(r"^\d{4}-\d{2}-\d{2}-", "", p.stem) for p in sorted(leads_dir.glob("*.md"))
    ]


def _scan_vault_folders() -> str:
    """Liest die aktuelle Ordnerstruktur des Vault dynamisch aus (Port aus
    _agent/heartbeat.py), statt einer hartcodierten Zielordner-Liste — damit
    neue Kunden und Unterordner (z.B. Meetings/Dokumente) automatisch bekannt sind."""
    settings = get_settings()
    lines = []
    for top in sorted(settings.vault_path.iterdir()):
        if top.name.startswith(".") or top.name.startswith("_") or not top.is_dir():
            continue
        if top.name in ("backend", "frontend", "Uni"):
            continue
        subs = sorted(p.name for p in top.iterdir() if p.is_dir() and not p.name.startswith("."))
        if subs:
            for sub in subs:
                subsubs = sorted(p.name for p in (top / sub).iterdir() if p.is_dir() and not p.name.startswith("."))
                if subsubs:
                    for ss in subsubs:
                        lines.append(f"- {top.name}/{sub}/{ss}")
                else:
                    lines.append(f"- {top.name}/{sub}")
        else:
            lines.append(f"- {top.name}")
    return "\n".join(lines)


_MEMORY_RULES_MAX_LINES = 60


def _load_memory_rules() -> str:
    """Liest NUR den '## REGEL'-Abschnitt aus memory.md - das sind echte,
    kuratierte Ablage-/Verhaltensregeln. '## PROZESS' wurde hier bis
    2026-08-14 mit eingelesen, ist aber tatsächlich ein unbegrenzt
    wachsendes, chronologisches Log einzelner Business-Fakten aus JEDER
    verarbeiteten Datei (memory.learn_from_file(), memory.py) - keine
    Ablageregeln. Das ließ den Abschnitt auf ~880 Zeilen/140KB anwachsen und
    blähte den classify()-Prompt so stark auf, dass die Klassifizierung
    mancher Dateien fehlschlug ("API-Klassifizierung fehlgeschlagen").
    Zusätzliche Zeilen-Obergrenze als Schutz, falls REGEL selbst irgendwann
    stark wächst - nimmt dann die neuesten Einträge (die stehen laut
    append_to_memory() direkt nach der Überschrift, also zuerst)."""
    settings = get_settings()
    if not settings.memory_path.exists():
        return ""
    lines = []
    in_section = False
    for line in settings.memory_path.read_text(encoding="utf-8").splitlines():
        if line.startswith("## REGEL"):
            in_section = True
        elif line.startswith("## ") and in_section:
            break
        elif in_section and line.startswith("- "):
            lines.append(line)
            if len(lines) >= _MEMORY_RULES_MAX_LINES:
                break
    return "\n".join(lines)


def classify(filepath: Path, content: str) -> dict | None:
    settings = get_settings()
    relative = filepath.relative_to(settings.inbox_dir)
    ordner_kontext = str(relative.parent) if str(relative.parent) != "." else ""
    memory_rules = _load_memory_rules()
    vault_struktur = _scan_vault_folders()

    prompt = f"""Du sortierst Dokumente für Prozessia GbR (KI-Agentur, Saarbrücken) ein.

Dateiname: {filepath.name}
Inbox-Unterordner: {ordner_kontext or "keiner"}
Inhalt (Auszug): {content}

{f"Gelernte Ablageregeln:{chr(10)}{memory_rules}{chr(10)}" if memory_rules else ""}
Aktuelle Vault-Struktur (alle existierenden Ordner):
{vault_struktur}

REGELN:
- Kundendokumente → Kunden/[Kundenname]/[Unterordner]
  Unterordner-Logik:
  · Vertraege/   → NDA, AVV, SLA, Bestellungen, Rechnungen, Auftragsbestätigungen
  · Angebote/    → Angebote, Kostenkalkulationen, Wartungsmodelle, Preislisten
  · Meetings/    → Besprechungen, Protokolle, Mitschriften, Meet-Aufzeichnungen
  · Praesentationen/ → Präsentationen, Pitches, Slides
  · Dokumente/   → alles andere (Specs, Fachkonzepte, Anleitungen, Projektpläne)
- Neuer Kunde erkannt → zielordner: "Kunden/[Firmenname]/Dokumente" (Ordner wird automatisch erstellt)
  ACHTUNG neuer_kunde=true nur bei einer ECHTEN neuen Kunden-/Interessenten-Beziehung
  für Prozessia selbst. KEIN neuer Kunde ist z.B.: ein externer Lieferant/
  Subunternehmer/IT-Dienstleister, der nur im Rahmen eines BESTEHENDEN
  Kundenprojekts erwähnt wird (der gehört zu diesem bestehenden Kunden, nicht
  in einen eigenen Ordner); eine private oder thematisch fremde Person/Inhalt
  ohne Geschäftsbezug zu Prozessia; ein Newsletter, eine Werbe-Mail oder
  generische Massenkommunikation, in der zufällig ein Name auftaucht. Im
  Zweifel eher zu Memos/ oder zum bereits bestehenden Kunden einordnen als
  einen neuen Kundenordner zu erfinden.
- Verträge die keinen Kunden zugeordnet werden können → Vertraege/
- Leads/Interessenten (noch kein Kunde, aber eine echte, mehrteilige
  Geschäftsbeziehung mit Meetings/Angeboten/Dokumenten): GENAUSO wie
  Kundendokumente behandeln → "Kunden/[Firmenname]/[Unterordner]", exakt
  dieselbe Unterordner-Logik wie oben. Es gibt KEINEN separaten
  "-Korrespondenz"-Ordner und KEINE Sonderbehandlung mehr für Leads mit
  mehreren Dokumenten - ob ein Vertrag schon unterschrieben ist oder nicht,
  entscheidet sich am Inhalt der Dokumente selbst (Status-Auswertung), nicht
  an einer eigenen Ordner-Hierarchie. Existiert in der Vault-Struktur unten
  schon ein Kunden/[Name]/-Ordner (auch bei anderer Schreibweise erkennbar
  ähnlich), IMMER diesen wiederverwenden statt einen neuen/parallelen Ordner
  anzulegen - das war die Hauptfehlerquelle für Dubletten.
  NUR ein ganz frisch erkannter Erstkontakt ohne inhaltliche Substanz (z.B.
  ein einzelner Kalendereintrag ohne Meeting-Notiz/Angebot/Dokument) bleibt
  eine flache Einzeldatei: "Leads/[Datum]-[Lead-Name].md", noch kein eigener
  Ordner. Sobald ein zweites Dokument zu diesem Lead dazukommt, wird daraus
  sofort ein "Kunden/[Firmenname]/"-Ordner (nicht "Leads/[Name]/") - dorthin
  auch die bereits bestehende Einzeldatei verschieben/zusammenführen.
- Finanzen → Finanzen/Rechnungen oder Finanzen/Angebote
- Marketing → Marketing/[passender Unterordner]
- Sales → Sales/[passender Unterordner]

Bestimme:
1. kategorie: Kunde/Lead/Produkt/Vertrag/Marketing/Sales/Finanzen/Memo
2. zusammenfassung: 2-3 Sätze was das Dokument enthält (bei Medien: kurze Beschreibung)
3. tags: 3-5 relevante Tags als JSON-Array
4. zielordner: exakter relativer Pfad (darf neu sein, wird erstellt)
5. neuer_kunde: true/false (ob ein bisher unbekannter Kunde erkannt wurde)
6. unsicher: true/false - ob du bei der Einordnung nicht sicher bist (z.B.
   Dokument könnte zu mehreren Kunden/Kategorien passen, Kundenname
   mehrdeutig, Inhalt unklar/widersprüchlich). Es läuft hier kein Mensch mit,
   der zwischendurch nachfragen könnte - im Zweifel lieber true setzen, damit
   der Fall sichtbar im Log auftaucht, statt es lautlos zu entscheiden.
7. unsicherheitsgrund: falls unsicher=true, ein kurzer deutscher Satz warum
   (sonst leerer String)

Antworte NUR als JSON, keine Erklärung."""

    try:
        # thinking explizit deaktiviert (Bug 2026-07-17): ohne das kann das Modell
        # einen Teil des knappen max_tokens-Budgets fürs Nachdenken verbrauchen und
        # die eigentliche JSON-Antwort abschneiden (stop_reason "max_tokens" mit
        # leerem/kaputtem Text-Block) - genau das ließ z.B. das TPG-Transkript als
        # "API-Klassifizierung fehlgeschlagen" durchfallen, obwohl der Aufruf an
        # sich funktioniert hätte.
        raw = complete_json(prompt, model=Models.SONNET, max_tokens=500).strip()
        raw = raw.replace("```json", "").replace("```", "").strip()
        result = json.loads(raw)
        if result.get("neuer_kunde"):
            ziel = settings.vault_path / result.get("zielordner", "Memos")
            kunde_root = ziel.parent if ziel.name in ("Vertraege", "Angebote", "Dokumente", "Meetings", "Praesentationen") else ziel
            for sub in ("Vertraege", "Angebote", "Dokumente", "Meetings"):
                (kunde_root / sub).mkdir(parents=True, exist_ok=True)
        return result
    except Exception:
        return None


def extract_meeting_structure(content: str) -> dict | None:
    """Extrahiert Teilnehmer/Kernpunkte/Zusagen/Nächste Schritte aus einem
    Meeting-Transkript. Nimmt den bereits mit großem max_chars extrahierten
    Volltext entgegen (siehe process_file()) statt die Datei ein zweites Mal
    einzulesen - vorher wurde hier pro Transkript eine eigene, zusätzliche
    Mistral-OCR-Extraktion ausgelöst, obwohl process_file() den Inhalt schon
    hatte (Sebastian, 2026-07-30)."""
    if not content:
        return None

    prompt = f"""Das ist das Transkript eines Kundengesprächs für Prozessia GbR.

Transkript:
{content}

Extrahiere:
1. teilnehmer: Namen der Gesprächsteilnehmer als JSON-Array (kurz, ohne Rollen/Firmen falls nicht eindeutig erkennbar)
2. kernpunkte: die wichtigsten besprochenen Punkte als JSON-Array kurzer deutscher Sätze (max 6)
3. zusagen: konkrete Zusagen, die im Gespräch gemacht wurden (von wem, was), als JSON-Array (leeres Array wenn keine)
4. naechste_schritte: offene Punkte / vereinbarte nächste Schritte als JSON-Array (leeres Array wenn keine)
5. entscheidungen: konkret getroffene Entscheidungen im Gespräch (z.B. "wir machen X statt Y", eine Richtung/ein Vorgehen wurde festgelegt) als JSON-Array kurzer deutscher Sätze - NUR echte Entscheidungen, keine offenen Fragen oder bloßen Diskussionspunkte (leeres Array wenn keine)
6. datum: falls im Transkript ein konkretes Datum des Gesprächs genannt wird (z.B. "heute ist der...", ein Datum im Header, eine Terminangabe für dieses Gespräch), im Format YYYY-MM-DD - sonst null

Antworte NUR als JSON, keine Erklärung. Format:
{{"teilnehmer": [...], "kernpunkte": [...], "zusagen": [...], "naechste_schritte": [...], "entscheidungen": [...], "datum": "YYYY-MM-DD" oder null}}"""

    try:
        raw = complete_json(prompt, model=Models.SONNET, max_tokens=800).strip()
        raw = raw.replace("```json", "").replace("```", "").strip()
        return json.loads(raw)
    except Exception:
        return None


def _resolve_datum(meeting_data: dict | None, filepath: Path) -> str:
    """Bestimmt das Ablage-Datum: bevorzugt ein im Transkript genanntes Datum,
    sonst die mtime der Original-Inbox-Datei (näher am echten Gesprächs-
    zeitpunkt als der Verarbeitungsmoment bei verzögertem Import), sonst
    'jetzt' als letzter Fallback. Behebt den Bug, dass verspätet verarbeitete
    Transkripte fälschlich das Verarbeitungsdatum statt des echten Meeting-
    Datums bekamen (Sebastian, 2026-07-18)."""
    heute = datetime.now().date()
    llm_datum = (meeting_data or {}).get("datum")
    if llm_datum:
        try:
            geparst = datetime.strptime(llm_datum, "%Y-%m-%d").date()
            if date(2020, 1, 1) <= geparst <= heute:
                return geparst.strftime("%Y-%m-%d")
        except (ValueError, TypeError):
            pass
    try:
        return datetime.fromtimestamp(filepath.stat().st_mtime).strftime("%Y-%m-%d")
    except OSError:
        return heute.strftime("%Y-%m-%d")


def _append_decisions_log(entscheidungen: list, quelle: str, datum: str) -> None:
    """Zentrales, vault-weites Entscheidungsprotokoll (Umsetzungsplan-Memo
    2026-07-16, Punkt B1 - "Decision Log" aus dem Everlast-AI-Company-Brain-
    Konzept). Ergänzung zu extract_meeting_structure(): jede dort erkannte
    Entscheidung landet zusätzlich hier gesammelt, damit sie nicht nur im
    einzelnen Meeting-Dokument, sondern durchsuchbar an einem Ort steht."""
    if not entscheidungen:
        return
    settings = get_settings()
    log_path = settings.agent_dir / "decisions.md"
    if not log_path.exists():
        log_path.parent.mkdir(parents=True, exist_ok=True)
        log_path.write_text("# Entscheidungsprotokoll\n\n", encoding="utf-8")
    with open(log_path, "a", encoding="utf-8") as f:
        for entscheidung in entscheidungen:
            f.write(f"- {datum} | {quelle} | {entscheidung}\n")


def is_meeting_transcript(filepath: Path, content: str, result: dict) -> bool:
    """Erkennt ein Gesprächstranskript unabhängig davon, in welche Kategorie das
    Modell die Datei einsortiert hat.

    Hintergrund (Sebastian, 2026-07-27): die Transkripte-Seite listet ausschließlich
    Notizen mit einem "Meetings"-Ordner im Pfad (files.py:list_meetings). Die
    Ablageregeln im classify()-Prompt kennen Meetings/ aber nur unter Kunden/ -
    Gespräche mit Interessenten landeten deshalb in Leads/[Lead]-Korrespondenz/
    und tauchten in der Übersicht nie auf, obwohl sie korrekt eingelesen waren
    (betraf u.a. Zillmer 21.07. und Seifert 23.07.). Zusätzlich blieben ihnen die
    Teilnehmer-/Kernpunkte-/Zusagen-Abschnitte verwehrt, weil auch die an dieser
    Ordnerbedingung hingen.

    Tag-Check bewusst als Whitelist exakter Tags, nicht als Substring-Match
    (Sebastian, 2026-07-28): ein Pitch-Deck über ein Transkriptions-/Voice-
    Produkt landete fälschlich hier, weil "transkript" auch in thematischen
    Tags wie "Transkription" oder "Live-Transkription" steckt (ist sogar ein
    Präfix davon) - das sagt nichts darüber aus, ob DIESES Dokument selbst ein
    Transkript ist."""
    haystack = f"{filepath.name}\n{content[:3000]}".lower()
    if any(marker in haystack for marker in TRANSCRIPT_MARKERS):
        return True
    exact_tags = {"transkript", "transkripte", "besprechungstranskript", "gesprächstranskript", "meeting-transkript"}
    return any(str(tag).strip().lower() in exact_tags for tag in result.get("tags", []))


def _extract_firma(zielordner_rel: str) -> str:
    """Leitet den Kunden-/Firmennamen aus dem Zielordner ab (Kunden/[Firma]/...),
    damit Notizen die Zugehörigkeit direkt im Frontmatter/Titel zeigen statt nur
    implizit über den Ordnerpfad (Sebastian, 2026-07-30: Transkript-Überschriften
    sollen klar zeigen, zu welcher Firma sie gehören).

    Das "-Korrespondenz"-Suffix unter Leads/ entstand nicht hier im
    classify()-Prompt (der 2026-08-11-Fix dort betraf nur diesen einen
    Schreibpfad), sondern in zwei von classify.py unabhängigen Stellen -
    email_indexer._write_lead_correspondence() und email_lead_service.
    consider_new_lead() -, die erst am 2026-08-19 entfernt wurden (Sebastian:
    "jeder Kunde hat ein Kundenordner, das reicht", inkl. Migration der
    bestehenden Leads/*-Korrespondenz/-Ordner nach Kunden/). Die re.sub hier
    bleibt als harmloser No-Op für den Fall, dass irgendwo doch noch ein
    Altordner mit dieser Endung existiert."""
    parts = Path(zielordner_rel).parts
    if len(parts) >= 2 and parts[0] == "Kunden":
        return parts[1]
    if len(parts) >= 2 and parts[0] == "Leads":
        return re.sub(r"-Korrespondenz$", "", parts[1])
    return ""


def process_file(filepath: Path) -> tuple[bool, str]:
    settings = get_settings()
    if filepath.suffix.lower() in SKIP_EXTENSIONS:
        return False, "Code-Datei übersprungen"

    if _is_trivial_image(filepath):
        _move_aside(filepath, settings.agent_dir / "trivial_images")
        return False, "Trivial-Bild übersprungen (nach _agent/trivial_images/ verschoben)"

    file_hash = None
    try:
        file_hash = _file_hash(filepath)
    except Exception:
        file_hash = None
    if file_hash:
        existing = _load_content_hashes().get(file_hash)
        if existing and (settings.vault_path / existing).exists():
            _move_aside(filepath, settings.inbox_dir / "_duplikate")
            return False, f"Duplikat von {existing}"

    # max_chars=200000 statt vormals 20000 (das wiederum vormals 3000 war): der
    # Volltext wird unten sowohl für die Meeting-Struktur als auch für die
    # gespeicherte Notiz wiederverwendet - bei 20000 Zeichen war bei jedem
    # etwas längeren Meeting (z.B. TopDown 14.08.2026: echtes Transkript
    # 35065 Zeichen) der "Vollständiger Inhalt"-Abschnitt immer noch mitten im
    # Gespräch abgeschnitten (Sebastian, 2026-08-17: "immer noch nicht
    # vollständig"). 200000 Zeichen (~35k Wörter) deckt auch mehrstündige
    # Meetings ab. Kostet keinen zusätzlichen API-Call: die Mistral-OCR-Antwort
    # wird nur lokal weniger stark gekürzt, PyPDF2/docx/etc. lesen ohnehin lokal.
    content = extract_text(filepath, max_chars=200000)
    if content is None:
        return False, "Extraktion fehlgeschlagen"

    # Leerer Extrakt separat behandeln (2026-08-11): vorher lief ein leerer
    # String weiter in classify(), das Modell bekam nichts zu klassifizieren
    # und die Datei landete mit der irreführenden Meldung
    # "API-Klassifizierung fehlgeschlagen" in _fehler/ - die eigentliche
    # Ursache (nichts extrahierbar) war aus dem Log nicht mehr ablesbar.
    # Bei PDFs ist das fast immer ein Scan ohne Textebene: dann greift die
    # Mistral-OCR-Vorstufe in extract_text() nur, wenn MISTRAL_API_KEY gesetzt
    # ist, sonst liefert der PyPDF2-Fallback leeren Text.
    if not content.strip():
        if filepath.suffix.lower() == ".pdf" and not settings.mistral_api_key:
            return False, "Kein Text extrahierbar (vermutlich Scan) - MISTRAL_API_KEY nicht gesetzt, OCR übersprungen"
        return False, f"Kein Text extrahierbar aus {filepath.suffix or 'Datei ohne Endung'}"

    result = classify(filepath, content[:3000])
    if not result:
        return False, "API-Klassifizierung fehlgeschlagen"

    # Transkripte immer in einen Meetings-Unterordner, egal ob Kunde, Lead oder
    # Sales - sonst fehlen sie in der Transkripte-Übersicht (siehe
    # is_meeting_transcript()). Hat das Modell bereits Meetings/ gewählt, bleibt
    # der Pfad unverändert.
    zielordner_rel = _sanitize_zielordner(result.get("zielordner", "Memos"))
    result["zielordner"] = zielordner_rel
    ist_transkript = is_meeting_transcript(filepath, content, result)
    if ist_transkript and "Meetings" not in Path(zielordner_rel).parts:
        zielordner_rel = f"{zielordner_rel}/Meetings"
        result["zielordner"] = zielordner_rel

    zielordner = settings.vault_path / zielordner_rel
    zielordner.mkdir(parents=True, exist_ok=True)
    firma = _extract_firma(zielordner_rel)

    # Original bleibt flach im Zielordner, die zugehörige Notiz liegt separat
    # unter <zielordner>/MD/ (Sebastian, 2026-08-14: .md-Notizen dürfen nie
    # mit Originaldokumenten vermischt in einem Ordner liegen - löst die
    # vorherige Regel vom 11.08. ab, die Original+Notiz gemeinsam in einem
    # eigenen, nach dem Dokument benannten Unterordner ablegte). Bilder/Zips
    # haben ohnehin keine Notiz und bleiben unverändert flach. Läuft ein
    # Zielordner trotzdem wieder voll, übernimmt reorganize_vault()
    # (min_dateien-Schwelle) die automatische Kategorie-Unterteilung.
    hat_notiz = filepath.suffix.lower() not in IMAGE_EXTS | {".zip"}
    if hat_notiz:
        (zielordner / "MD").mkdir(exist_ok=True)

    meeting_data = extract_meeting_structure(content) if ist_transkript else None
    datum = _resolve_datum(meeting_data, filepath)

    if hat_notiz:
        tags = result.get("tags", [])
        tags_str = "\n".join(f"  - {t}" for t in tags)
        notiz_path = zielordner / "MD" / f"{datum}-{filepath.stem}.md"

        meeting_sections = ""
        if ist_transkript:
            if meeting_data:
                def _liste(items):
                    return "\n".join(f"- {i}" for i in items) if items else "- (keine)"
                meeting_sections = f"""
## Teilnehmer
{_liste(meeting_data.get("teilnehmer", []))}

## Kernpunkte
{_liste(meeting_data.get("kernpunkte", []))}

## Zusagen
{_liste(meeting_data.get("zusagen", []))}

## Nächste Schritte
{_liste(meeting_data.get("naechste_schritte", []))}

## Entscheidungen
{_liste(meeting_data.get("entscheidungen", []))}
"""
                _append_decisions_log(meeting_data.get("entscheidungen", []), filepath.name, datum)

        # Firma/Teilnehmer explizit im Frontmatter und in der Überschrift, damit
        # die Zugehörigkeit einer Notiz (v.a. bei Transkripten) auf einen Blick
        # erkennbar ist, statt sich nur aus dem Ordnerpfad zu ergeben (Sebastian,
        # 2026-07-30).
        frontmatter_extra = f"firma: {firma}\n" if firma else ""
        if meeting_data and meeting_data.get("teilnehmer"):
            frontmatter_extra += f'teilnehmer: "{", ".join(meeting_data["teilnehmer"])}"\n'
        titel_praefix = f"[{firma}] " if firma and ist_transkript else ""

        notiz_path.write_text(
            f"""---
tags:
{tags_str}
quelle: {filepath.name}
datum: {datum}
kategorie: {result.get("kategorie", "")}
{frontmatter_extra}---

# {titel_praefix}{filepath.stem}

## Zusammenfassung
{result.get("zusammenfassung", "")}
{meeting_sections}
## Vollständiger Inhalt
{content}
""",
            encoding="utf-8",
        )

    safe_name = filepath.name.replace("@", "-at-")
    ziel_original = zielordner / safe_name
    ziel = ziel_original if not ziel_original.exists() else zielordner / f"{datum}-{safe_name}"
    try:
        shutil.move(str(filepath), str(ziel))
    except Exception:
        try:
            shutil.copy2(str(filepath), str(ziel))
            filepath.unlink(missing_ok=True)
        except Exception:
            pass

    # Dedup-Hash für künftige Duplikat-Erkennung hinterlegen (siehe
    # _load_content_hashes() oben) - erst NACH dem erfolgreichen Move, damit
    # der referenzierte Pfad auch wirklich existiert.
    if file_hash:
        try:
            hashes = _load_content_hashes()
            hashes[file_hash] = str(ziel.relative_to(settings.vault_path))
            _save_content_hashes(hashes)
        except Exception:
            pass

    unsicher = bool(result.get("unsicher"))
    log_path = settings.agent_dir / "logs" / "inbox_log.md"
    log_path.parent.mkdir(parents=True, exist_ok=True)
    with open(log_path, "a", encoding="utf-8") as f:
        marker = "⚠️ UNSICHER | " if unsicher else ""
        grund = f" | Grund: {result.get('unsicherheitsgrund', '')}" if unsicher else ""
        f.write(
            f"- {datetime.now().strftime('%Y-%m-%d %H:%M')} | {marker}{filepath.name} → "
            f"{result.get('zielordner')} | {result.get('zusammenfassung', '')[:80]}{grund}\n"
        )

    info = result.get("zielordner", "")
    if unsicher:
        info = "⚠️ UNSICHER: " + info
    return True, info


def _ordnername(name: str) -> str:
    """Sanitisiert einen freien Namen (Dokument-Stamm oder LLM-Kategorie) zu
    einem gültigen Ordnernamen: Pfadtrenner/führende-trailing Whitespaces raus,
    Länge gedeckelt, damit z.B. komplette E-Mail-Betreffzeilen oder etwas zu
    ausführliche Kategorienamen nicht an Dateisystem-Limits scheitern."""
    name = name.strip().replace("/", "-").replace("\\", "-")
    return name[:80] if len(name) > 80 else name


VERTRAEGE_KATEGORIEN = [
    "NDA", "Dienstleistungsvertraege", "Vertriebsvereinbarungen",
    "AVV_SLA", "Gesellschaftsrecht", "Sonstige",
]

# Top-Level-Ordner, die reorganize_vault() als Inhaltsordner behandelt -
# bewusst eine Allowlist statt einer Blockliste, damit backend/, frontend/,
# dev-agent/, mcp-vnc/, services/, _agent/, _inbox/ sicher nie angefasst
# werden, ohne dass man jeden davon einzeln ausschließen muss.
REORG_TOP_LEVEL = {"Kunden", "Leads", "Sales", "Vertraege", "Finanzen", "Marketing", "Produkte", "Memos"}


_REORG_GRACE_SECONDS = 300  # 5 Minuten


def _kuerzlich_geaendert(pfad: Path) -> bool:
    """Ob `pfad` innerhalb der letzten _REORG_GRACE_SECONDS veraendert wurde -
    Sicherheitsabstand gegen Wettlaeufe mit externen Live-Sync-Tools (z.B.
    Syncthing), die dieselbe Datei gerade noch schreiben/verschieben koennten
    (Sebastian, 2026-08-14: beim Erst-Sync von ~1300 bisher VPS-unsichtbaren
    Dateien griff group_by_category() zeitgleich auf frisch von Syncthing
    geschriebene Dateien zu - 3 Dateien wurden kurzzeitig geloescht, per
    Syncthing-Versionierung wiederhergestellt, kein dauerhafter Verlust, aber
    ein echter Wettlauf). Bei fehlendem/nicht lesbarem stat() konservativ
    true (lieber diesen Zyklus ueberspringen als riskieren) - der naechste
    Lauf (reorganize_vault laeuft periodisch) holt es dann nach."""
    try:
        return (time.time() - pfad.stat().st_mtime) < _REORG_GRACE_SECONDS
    except OSError:
        return True


def unwrap_legacy_pairing_folders(parent: Path) -> dict:
    """Löst die bis 2026-08-14 aktive Struktur "ein Unterordner pro Dokument
    mit Original+Notiz zusammen" (siehe process_file() vor diesem Datum)
    wieder auf: ein direkter Unterordner von `parent`, der KEINE eigenen
    Unterordner, aber genau eine Nicht-MD-Datei und mindestens eine
    .md-Datei enthält, gilt als solcher Wrapper. Original wandert nach
    `parent/`, die .md(s) nach `parent/MD/`, der jetzt leere Wrapper wird
    entfernt. Muss vor separate_md_files() laufen, sonst hält dessen
    "enthält Ordner MD schon?"-Logik einen 1-Datei-Wrapper fälschlich für
    bereits korrekt getrennt. Kollisionsfälle (Zielname existiert schon)
    werden übersprungen statt überschrieben - idempotent, ein zweiter Lauf
    auf einem bereits aufgelösten Ordner findet nichts mehr zu tun."""
    if not parent.is_dir():
        return {"ordner": str(parent), "aufgeloest": 0, "moves": []}

    moves: list[tuple[Path, Path]] = []
    aufgeloest = 0
    for sub in [p for p in parent.iterdir() if p.is_dir() and p.name != "MD" and not p.name.startswith(".")]:
        kinder = [f for f in sub.iterdir() if not f.name.startswith(".")]
        if any(f.is_dir() for f in kinder):
            continue
        md_dateien = [f for f in kinder if f.suffix.lower() == ".md"]
        andere = [f for f in kinder if f.suffix.lower() != ".md"]
        if len(andere) != 1 or not md_dateien:
            continue
        if any(_kuerzlich_geaendert(f) for f in kinder):
            continue  # noch frisch (z.B. von einem Live-Sync-Tool) - naechster Zyklus

        original = andere[0]
        ziel_original = parent / original.name
        md_ziel_ordner = parent / "MD"
        if ziel_original.exists() or any((md_ziel_ordner / md.name).exists() for md in md_dateien):
            continue  # Kollision - lieber liegen lassen als überschreiben

        md_ziel_ordner.mkdir(exist_ok=True)
        try:
            shutil.move(str(original), str(ziel_original))
            for md in md_dateien:
                zielpfad = md_ziel_ordner / md.name
                shutil.move(str(md), str(zielpfad))
                moves.append((md, zielpfad))
        except (FileNotFoundError, shutil.Error):
            # Quelle ist zwischen Auflisten und Verschieben verschwunden -
            # z.B. weil derselbe Ordner gleichzeitig vom periodischen
            # Hintergrund-Job (jobs.py:vault_reorganize_loop) bearbeitet
            # wurde (Sebastian, 2026-08-14: bei einem manuell angestoßenen
            # Testlauf parallel zum Loop beobachtet). Einzelnen Wrapper
            # überspringen statt den ganzen Durchlauf abzubrechen.
            continue
        try:
            sub.rmdir()
        except OSError:
            pass
        aufgeloest += 1

    return {"ordner": str(parent), "aufgeloest": aufgeloest, "moves": moves}


def separate_md_files(folder: Path) -> dict:
    """Trennt lose .md-Notizen von allem anderen in `folder`: sobald der
    Ordner direkt sowohl mindestens eine .md-Datei als auch mindestens eine
    andere Datei oder einen Unterordner enthält, wandern alle direkten
    .md-Kinder nach `folder/MD/` (Sebastian, 2026-08-14: ".md-Dateien mit
    anderen vermischt - das geht nicht"). No-op wenn `folder` selbst "MD"
    heißt oder keine .md-Datei direkt enthält - macht die Funktion idempotent,
    ein zweiter Lauf auf einem bereits getrennten Ordner findet nichts mehr
    zu tun (die .md-Dateien liegen dann schon in folder/MD/, nicht mehr
    direkt in folder)."""
    if not folder.is_dir() or folder.name == "MD":
        return {"ordner": str(folder), "verschoben": 0, "moves": []}

    direkte = [f for f in folder.iterdir() if not f.name.startswith(".")]
    md_dateien = [f for f in direkte if f.is_file() and f.suffix.lower() == ".md"]
    andere = [f for f in direkte if f not in md_dateien]

    if not md_dateien or not andere:
        return {"ordner": str(folder), "verschoben": 0, "moves": []}

    stabile_md_dateien = [f for f in md_dateien if not _kuerzlich_geaendert(f)]
    if not stabile_md_dateien:
        return {"ordner": str(folder), "verschoben": 0, "moves": []}

    ziel = folder / "MD"
    ziel.mkdir(exist_ok=True)
    moves: list[tuple[Path, Path]] = []
    for md in stabile_md_dateien:
        zielpfad = ziel / md.name
        if zielpfad.exists():
            continue
        try:
            shutil.move(str(md), str(zielpfad))
        except (FileNotFoundError, shutil.Error):
            continue  # siehe Kommentar in unwrap_legacy_pairing_folders() - Race mit parallelem Lauf
        moves.append((md, zielpfad))
    return {"ordner": str(folder), "verschoben": len(moves), "moves": moves}


_ASSET_ORDNER_SUFFIXE = ("_files", "_Dateien")


def _sieht_aus_wie_echte_dokumente(ordner: Path, dateien: list[Path]) -> bool:
    """Filtert Cache-/Asset-Ordner aus der automatischen Kategorisierung
    heraus (Sebastian, 2026-08-14: beobachtet bei einem im Browser
    gespeicherten LinkedIn-Profil - "Seite_files/" mit hunderten hash-
    benannten Dateien ohne Endung; group_by_category() erfand dafür
    bedeutungslose Kategorien wie "Gruppe_A"/"Zahlen_1bis3", weil solche
    Ordner denselben Schwellwert wie echte Dokumentordner überschreiten).
    Zwei Signale: (1) der Ordner selbst oder ein Vorfahre bis REORG_TOP_LEVEL
    endet auf ein bekanntes Browser-"gespeicherte-Seite"-Suffix, (2) weniger
    als die Hälfte der Dateien hat überhaupt eine Dateiendung - echte
    Dokumente (PDF, DOCX, Bilder, ...) haben praktisch immer eine, rohe
    Cache-/Hash-Dateien meist nicht."""
    for teil in ordner.parts:
        if teil.endswith(_ASSET_ORDNER_SUFFIXE):
            return False
    mit_endung = sum(1 for f in dateien if f.suffix)
    return mit_endung >= len(dateien) / 2


def group_by_category(folder: Path, dateien: list[Path], vokabular: list[str] | None = None) -> dict:
    """Gruppiert lose Originaldateien in `folder` per LLM in Kategorie-
    Unterordner - mit festem Vokabular (z.B. Vertragstypen für Vertraege/,
    siehe VERTRAEGE_KATEGORIEN) oder, wenn keins übergeben wird, schlägt das
    Modell selbst 3-6 knappe Kategorien für genau diesen Dateisatz vor
    (Sebastian, 2026-08-14: "ab einer gewissen Zahl an Dateien müssen
    automatisch und intelligent Unterordner angelegt werden"). Nutzt nur die
    Dateinamen (kein Inhaltsauswertung nötig, hält den Aufruf günstig und
    schnell für bis zu einige hundert Dateien in einem Aufruf).

    Verschiebt NUR die übergebenen Originaldateien - .md-Notizen liegen nach
    separate_md_files() ohnehin schon getrennt in folder/MD/ und bleiben dort
    unverändert liegen; die Zuordnung zu einem Original bleibt über das
    "quelle:"-Frontmatter-Feld in der Notiz nachvollziehbar, auch ohne
    gemeinsamen Ordner."""
    if not dateien:
        return {"ordner": str(folder), "verschoben": 0, "moves": []}

    namen = [f.name for f in dateien]
    vokabular_hinweis = (
        f"Nutze GENAU diese Kategorien: {', '.join(vokabular)}"
        if vokabular else
        "Schlage selbst 3-6 knappe, treffende deutsche Kategorienamen vor "
        "(kurz, ohne Sonderzeichen/Leerzeichen, z.B. 'Rechnungen' statt "
        "'Rechnungen und Belege')."
    )
    prompt = f"""Ordne jede dieser Dateien einer Kategorie zu, um einen vollgelaufenen
Ordner ({folder.name}, {len(namen)} Dateien) in Unterordner aufzuteilen.

Dateinamen:
{chr(10).join(f"- {n}" for n in namen)}

{vokabular_hinweis}

Antworte NUR als JSON-Objekt: {{"dateiname": "kategorie", ...}} - für JEDE
Datei aus der Liste oben exakt ein Eintrag, "kategorie" als kurzer
Ordnername (keine Sonderzeichen außer _)."""

    try:
        raw = complete_json(prompt, model=Models.SONNET, max_tokens=2000).strip()
        raw = raw.replace("```json", "").replace("```", "").strip()
        zuordnung = json.loads(raw)
    except Exception:
        logger.exception("group_by_category: Klassifizierung fehlgeschlagen für %s", folder)
        return {"ordner": str(folder), "verschoben": 0, "moves": []}

    moves: list[tuple[Path, Path]] = []
    for f in dateien:
        kategorie = zuordnung.get(f.name)
        if not kategorie or not isinstance(kategorie, str):
            continue
        kategorie_ordner = folder / _ordnername(kategorie)
        kategorie_ordner.mkdir(exist_ok=True)
        ziel = kategorie_ordner / f.name
        zaehler = 1
        while ziel.exists():
            ziel = kategorie_ordner / f"{f.stem}({zaehler}){f.suffix}"
            zaehler += 1
        try:
            shutil.move(str(f), str(ziel))
        except (FileNotFoundError, shutil.Error):
            continue  # siehe Kommentar in unwrap_legacy_pairing_folders() - Race mit parallelem Lauf
        moves.append((f, ziel))

    return {"ordner": str(folder), "verschoben": len(moves), "moves": moves}


def reorganize_vault(min_dateien: int = 15) -> list[dict]:
    """Räumt den gesamten Vault auf, in drei Schritten über eine feste
    Allowlist von Top-Level-Inhaltsordnern (REORG_TOP_LEVEL, rekursiv in
    jeden Unterordner):

    1. unwrap_legacy_pairing_folders() - löst alte Dokument+Notiz-Wrapper-
       Ordner von vor dem 14.08. auf.
    2. separate_md_files() - trennt .md-Notizen von Originaldokumenten in
       jedem Ordner, der beides direkt gemischt enthält.
    3. group_by_category() - für Ordner, die danach noch >= min_dateien lose
       Originaldateien direkt enthalten, per LLM in Kategorie-Unterordner
       aufteilen (Vertraege/ mit festem Vertragstyp-Vokabular, sonst lässt
       das Modell selbst passende Kategorien vorschlagen).

    Reihenfolge ist bewusst: erst 1+2 über ALLE Kandidaten (bottom-up, damit
    ein aufgelöster Wrapper nicht seinen jetzt größeren Elternordner verpasst),
    danach erst 3 - sonst würde Schritt 3 noch die alten Wrapper-Ordner statt
    der echten losen Dateien zählen. Wird sowohl einmalig manuell aufgerufen
    als auch periodisch vom Hintergrund-Job (jobs.py:vault_reorganize_loop) -
    jede Teilfunktion ist für sich idempotent, siehe deren Docstrings. Zieht
    danach den RAG-Index nach (alte Pfade raus, neue Pfade rein) für alle
    verschobenen .md-Notizen - reine Dokument-Umsortierungen aus Schritt 3
    berühren den RAG-Index nicht, der indiziert nur .md-Dateien (rag.py)."""
    settings = get_settings()
    ergebnisse: list[dict] = []
    alte_pfade: set[str] = set()

    def sammle_dirs(basis: Path) -> list[Path]:
        dirs = [basis]
        for p in sorted(basis.rglob("*")):
            if p.is_dir() and not any(teil.startswith(".") for teil in p.relative_to(basis).parts):
                dirs.append(p)
        return dirs

    kandidaten: list[Path] = []
    for name in sorted(REORG_TOP_LEVEL):
        basis_dir = settings.vault_path / name
        if basis_dir.exists() and basis_dir.is_dir():
            kandidaten.extend(sammle_dirs(basis_dir))
    # Tiefste Ordner zuerst, siehe Docstring oben.
    kandidaten.sort(key=lambda p: len(p.parts), reverse=True)

    for ordner in kandidaten:
        if not ordner.exists():
            continue
        r = unwrap_legacy_pairing_folders(ordner)
        if r["aufgeloest"]:
            ergebnisse.append({"ordner": str(ordner), "typ": "unwrap", **r})
            for alt, _neu in r["moves"]:
                alte_pfade.add(str(alt.relative_to(settings.vault_path)))

    for ordner in kandidaten:
        if not ordner.exists() or ordner.name == "MD":
            continue
        r = separate_md_files(ordner)
        if r["verschoben"]:
            ergebnisse.append({"ordner": str(ordner), "typ": "md_trennung", **r})
            for alt, _neu in r["moves"]:
                alte_pfade.add(str(alt.relative_to(settings.vault_path)))

    for ordner in kandidaten:
        if not ordner.exists() or ordner.name == "MD":
            continue
        lose = [f for f in ordner.iterdir() if f.is_file() and not f.name.startswith(".")]
        stabile = [f for f in lose if not _kuerzlich_geaendert(f)]
        if len(stabile) < min_dateien or not _sieht_aus_wie_echte_dokumente(ordner, stabile):
            continue
        rel = ordner.relative_to(settings.vault_path)
        vokabular = VERTRAEGE_KATEGORIEN if rel.parts and rel.parts[0] == "Vertraege" else None
        r = group_by_category(ordner, stabile, vokabular)
        if r["verschoben"]:
            ergebnisse.append({"ordner": str(ordner), "typ": "kategorisierung", **r})

    if alte_pfade:
        try:
            from app.services import rag
            if rag.is_loaded():
                rag.remove_paths(alte_pfade)
                rag.reindex_new_files()
        except Exception:
            logger.exception("RAG-Index-Abgleich nach reorganize_vault() fehlgeschlagen")
    return ergebnisse


def _archive_source(datei: Path, ziel_ordner: Path) -> None:
    """Verschiebt eine Quelldatei nach _verarbeitet/, falls sie noch in
    _inbox/ liegt. Normalerweise hat process_file() sie beim erfolgreichen
    Einsortieren schon an ihren endgültigen Ort im Vault verschoben - das
    ist der Erfolgsfall, kein Fehler (Fix 2026-08-14: vorher versuchte diese
    Funktion trotzdem IMMER den Move und loggte bei JEDER erfolgreich
    verarbeiteten Datei eine irreführende "konnte nicht verschieben"-Warnung,
    weil die Quelle da schon lange weg war - live reproduziert, kein
    Datenverlust, nur Log-Rauschen). Nur falls process_file() ausnahmsweise
    weder verschieben noch kopieren konnte (seltener Doppel-Fehlerfall aus
    process_file()s eigenem try/except), liegt die Datei noch in _inbox/ -
    dann greift dieser Schritt als Sicherheitsnetz.

    Bei Namenskollision wird durchnummeriert, damit nie eine Datei überschrieben
    wird (gleiche Anhänge kommen über Mail-Weiterleitungen mehrfach rein)."""
    if not datei.exists():
        return
    ziel = ziel_ordner / datei.name
    zaehler = 1
    while ziel.exists():
        ziel = ziel_ordner / f"{datei.stem}({zaehler}){datei.suffix}"
        zaehler += 1
    try:
        shutil.move(str(datei), str(ziel))
    except Exception:
        logger.warning("Konnte %s nicht nach _verarbeitet/ verschieben", datei.name, exc_info=True)


def run_inbox() -> dict:
    """Verarbeitet alle neuen Dateien in _inbox/. Entspricht heartbeat.py:run()."""
    settings = get_settings()
    fehler_path = settings.inbox_dir / "_fehler"
    verarbeitet_path = settings.inbox_dir / "_verarbeitet"
    daily_path = settings.agent_dir / "daily"
    fehler_path.mkdir(parents=True, exist_ok=True)
    verarbeitet_path.mkdir(parents=True, exist_ok=True)
    daily_path.mkdir(parents=True, exist_ok=True)

    dateien = [
        f for f in settings.inbox_dir.rglob("*")
        if f.is_file()
        and not f.name.startswith(".")
        and f.suffix.lower() not in SKIP_EXTENSIONS
        and f.name not in SKIP_NAMES
        and not any(f.name.startswith(p) for p in SKIP_PREFIXES)
        and "node_modules" not in str(f)
        and "_fehler" not in str(f)
        and "_verarbeitet" not in str(f)
    ]

    cache = _load_cache()
    dateien = [d for d in dateien if str(d) not in cache]

    if not dateien:
        return {"processed": 0, "errors": 0, "output": "Inbox leer (alle Dateien bereits verarbeitet)."}

    processed = 0
    errors = 0
    log_lines = []
    for datei in dateien:
        try:
            ok, info = process_file(datei)
        except Exception as e:
            ok, info = False, f"Unerwarteter Fehler: {e}"
        if ok:
            cache.add(str(datei))
            _save_cache(cache)
            processed += 1
            log_lines.append(f"{datei.name} -> {info}")
            # Quelldatei nach _verarbeitet/ wegräumen (2026-08-11). Vorher blieb
            # sie nach erfolgreicher Einsortierung in _inbox/ liegen und wurde
            # nur über den Cache stillgelegt. Dadurch war am Ordner nicht mehr
            # ablesbar, was noch offen ist: zuletzt lagen 141 Dateien in _inbox,
            # 95 davon längst verarbeitet. Genau das Karteileichen-Problem, das
            # im Kommentar in background/jobs.py für den 24.07. beschrieben ist.
            # Bewusst verschieben statt löschen - das Original bleibt erhalten
            # (Vault-Regel: nie ohne Sebastians Bestätigung löschen).
            _archive_source(datei, verarbeitet_path)
        else:
            errors += 1
            log_lines.append(f"{datei.name}: FEHLER {info}")
            # Diese drei Fälle hat process_file() bereits selbst an ihren
            # endgültigen Platz verschoben (Code-Datei bleibt liegen, Dublette
            # -> _inbox/_duplikate/, Trivial-Bild -> _agent/trivial_images/) -
            # ein zusätzlicher Move nach _fehler/ würde das nur rückgängig
            # machen bzw. eine schon verschobene (nicht mehr existierende)
            # Datei erneut anfassen.
            bereits_verschoben = (
                info == "Code-Datei übersprungen"
                or info.startswith("Duplikat von ")
                or info.startswith("Trivial-Bild übersprungen")
            )
            if not bereits_verschoben:
                try:
                    shutil.move(str(datei), str(fehler_path / datei.name))
                except Exception:
                    pass

    daily_file = daily_path / f"{datetime.now().strftime('%Y-%m-%d')}.md"
    with open(daily_file, "a", encoding="utf-8") as f:
        f.write(f"\n## Inbox-Verarbeitung {datetime.now().strftime('%H:%M')}\n")
        f.write(f"- Verarbeitet: {processed}\n")
        f.write(f"- Fehler: {errors}\n")
        # Grund pro Datei mitschreiben (2026-08-11): log_lines wurde bisher nur
        # als Rückgabewert weitergereicht und nirgends persistiert - im Log
        # standen ausschließlich Zählerstände. Bei den Läufen mit 51-53 Fehlern
        # ließ sich hinterher nicht mehr feststellen, woran die Dateien
        # gescheitert sind; nachträglich getestet waren 64 der 65 PDFs in
        # _fehler/ problemlos lesbar, der Fehler lag also nicht bei der
        # Extraktion. Ohne diese Zeilen bleibt so etwas unauffindbar.
        for line in log_lines:
            f.write(f"  - {line}\n")

    return {"processed": processed, "errors": errors, "output": "\n".join(log_lines)}
