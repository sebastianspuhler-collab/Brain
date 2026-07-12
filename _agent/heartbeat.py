import os
import json
import shutil
from datetime import datetime
from pathlib import Path
import anthropic

VAULT_PATH = Path(__file__).parent.parent
INBOX_PATH = VAULT_PATH / "_inbox"
FEHLER_PATH = VAULT_PATH / "_inbox" / "_fehler"
LOGS_PATH = VAULT_PATH / "_agent" / "logs"
DAILY_PATH = VAULT_PATH / "_agent" / "daily"
API_KEY = os.environ.get("ANTHROPIC_API_KEY")

SKIP_EXTENSIONS = {".js", ".ts", ".map", ".mjs", ".jsx", ".tsx", ".css",
                   ".yml", ".yaml", ".eslintrc", ".nycrc", ".npmignore",
                   ".enc", ".lock", ".editorconfig", ".orig"}

FEHLER_PATH.mkdir(parents=True, exist_ok=True)
LOGS_PATH.mkdir(parents=True, exist_ok=True)
DAILY_PATH.mkdir(parents=True, exist_ok=True)
CACHE_PATH = LOGS_PATH / "processed_cache.json"

def load_cache():
    if CACHE_PATH.exists():
        try:
            return set(json.loads(CACHE_PATH.read_text()))
        except Exception:
            return set()
    return set()

def save_cache(cache):
    CACHE_PATH.write_text(json.dumps(list(cache)))

def extract_text(filepath, max_chars=3000):
    suffix = filepath.suffix.lower()
    try:
        if suffix == ".pdf":
            import PyPDF2
            with open(filepath, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                return " ".join(page.extract_text() or "" for page in reader.pages)[:max_chars]
        elif suffix == ".docx":
            from docx import Document
            doc = Document(filepath)
            return " ".join(p.text for p in doc.paragraphs)[:max_chars]
        elif suffix in [".xlsx", ".xls"]:
            import openpyxl
            wb = openpyxl.load_workbook(filepath, read_only=True)
            text = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    text.append(" ".join(str(c) for c in row if c))
            return " ".join(text)[:max_chars]
        elif suffix in [".pptx", ".ppt"]:
            from pptx import Presentation
            prs = Presentation(filepath)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return " ".join(text)[:max_chars]
        elif suffix == ".html":
            from html.parser import HTMLParser
            class TextExtractor(HTMLParser):
                def __init__(self):
                    super().__init__()
                    self.parts = []
                def handle_data(self, data):
                    self.parts.append(data)
            p = TextExtractor()
            p.feed(filepath.read_text(errors="ignore"))
            return " ".join(p.parts)[:max_chars]
        elif suffix == ".json":
            raw = filepath.read_text(errors="ignore")[:max_chars]
            return f"[JSON-Datei] {raw}"
        elif suffix in [".txt", ".md", ".markdown", ".csv"]:
            return filepath.read_text(errors="ignore")[:max_chars]
        elif suffix in [".png", ".jpg", ".jpeg", ".gif", ".svg", ".webp"]:
            return f"[Bilddatei, kein Text extrahierbar]"
        elif suffix in [".mp4", ".mov", ".avi", ".mkv", ".mp3"]:
            return f"[Mediendatei, kein Text extrahierbar]"
        elif suffix == ".zip":
            return f"[ZIP-Archiv]"
        else:
            return f"[Unbekanntes Format: {suffix}]"
    except Exception:
        return None

def load_memory_rules():
    memory_path = VAULT_PATH / "_agent" / "memory.md"
    if memory_path.exists():
        text = memory_path.read_text()
        # extract REGEL and PROZESS sections only
        lines = []
        in_section = False
        for line in text.splitlines():
            if line.startswith("## REGEL") or line.startswith("## PROZESS"):
                in_section = True
            elif line.startswith("## ") and in_section:
                in_section = False
            elif in_section and line.startswith("- "):
                lines.append(line)
        return "\n".join(lines) if lines else ""
    return ""


def _scan_vault_folders() -> str:
    """Liest die aktuelle Ordnerstruktur des Vault dynamisch aus."""
    lines = []
    for top in sorted(VAULT_PATH.iterdir()):
        if top.name.startswith(".") or top.name.startswith("_") or not top.is_dir():
            continue
        if top.name in ["backend", "frontend", "Uni"]:
            continue
        subs = sorted(p.name for p in top.iterdir() if p.is_dir() and not p.name.startswith("."))
        if subs:
            for sub in subs:
                subsubs = sorted(p.name for p in (top / sub).iterdir() if p.is_dir() and not p.name.startswith("."))
                if subsubs:
                    for ss in subsubs:
                        lines.append(f"- {top.name}/{sub}/{ss}")
                else:
                    lines.append(f"- {top.name}/{sub}")
        else:
            lines.append(f"- {top.name}")
    return "\n".join(lines)


def classify(filepath, content):
    client = anthropic.Anthropic(api_key=API_KEY)

    relative = filepath.relative_to(INBOX_PATH)
    ordner_kontext = str(relative.parent) if str(relative.parent) != "." else ""
    memory_rules = load_memory_rules()
    vault_struktur = _scan_vault_folders()

    prompt = f"""Du sortierst Dokumente für Prozessia GbR (KI-Agentur, Saarbrücken) ein.

Dateiname: {filepath.name}
Inbox-Unterordner: {ordner_kontext or "keiner"}
Inhalt (Auszug): {content}

{f"Gelernte Ablageregeln:{chr(10)}{memory_rules}{chr(10)}" if memory_rules else ""}
Aktuelle Vault-Struktur (alle existierenden Ordner):
{vault_struktur}

REGELN:
- Kundendokumente → Kunden/[Kundenname]/[Unterordner]
  Unterordner-Logik:
  · Vertraege/   → NDA, AVV, SLA, Bestellungen, Rechnungen, Auftragsbestätigungen
  · Angebote/    → Angebote, Kostenkalkulationen, Wartungsmodelle, Preislisten
  · Meetings/    → Besprechungen, Protokolle, Mitschriften, Meet-Aufzeichnungen
  · Praesentationen/ → Präsentationen, Pitches, Slides
  · Dokumente/   → alles andere (Specs, Fachkonzepte, Anleitungen, Projektpläne)
- Neuer Kunde erkannt → zielordner: "Kunden/[Firmenname]/Dokumente" (Ordner wird automatisch erstellt)
- Verträge die keinen Kunden zugeordnet werden können → Vertraege/
- Leads → Leads/
- Finanzen → Finanzen/Rechnungen oder Finanzen/Angebote
- Marketing → Marketing/[passender Unterordner]
- Sales → Sales/[passender Unterordner]

Bestimme:
1. kategorie: Kunde/Lead/Produkt/Vertrag/Marketing/Sales/Finanzen/Memo
2. zusammenfassung: 2-3 Sätze Inhalt
3. tags: 3-5 Tags als JSON-Array
4. zielordner: exakter relativer Pfad (darf neu sein, wird erstellt)
5. neuer_kunde: true/false (ob ein bisher unbekannter Kunde erkannt wurde)

Antworte NUR als JSON, keine Erklärung."""

    try:
        resp = client.messages.create(
            model="claude-sonnet-5",
            max_tokens=500,
            messages=[{"role": "user", "content": prompt}]
        )
        raw = resp.content[0].text.strip()
        raw = raw.replace("```json", "").replace("```", "").strip()
        result = json.loads(raw)
        # Wenn neuer Kunde: Standard-Unterordner anlegen
        if result.get("neuer_kunde"):
            ziel = VAULT_PATH / result.get("zielordner", "Memos")
            kunde_root = ziel.parent if ziel.name in ["Vertraege","Angebote","Dokumente","Meetings","Praesentationen"] else ziel
            for sub in ["Vertraege", "Angebote", "Dokumente", "Meetings"]:
                (kunde_root / sub).mkdir(parents=True, exist_ok=True)
        return result
    except Exception:
        return None


def extract_meeting_structure(filepath):
    """Extrahiert Teilnehmer/Kernpunkte/Zusagen/Nächste Schritte aus einem
    Meeting-Transkript. Liest die Datei mit größerem max_chars erneut ein,
    damit auch spät im Gespräch genannte Zusagen/nächste Schritte erfasst
    werden (die reguläre Klassifizierung kappt bei 3000 Zeichen)."""
    content = extract_text(filepath, max_chars=20000)
    if not content:
        return None

    client = anthropic.Anthropic(api_key=API_KEY)
    prompt = f"""Das ist das Transkript eines Kundengesprächs für Prozessia GbR.

Transkript:
{content}

Extrahiere:
1. teilnehmer: Namen der Gesprächsteilnehmer als JSON-Array (kurz, ohne Rollen/Firmen falls nicht eindeutig erkennbar)
2. kernpunkte: die wichtigsten besprochenen Punkte als JSON-Array kurzer deutscher Sätze (max 6)
3. zusagen: konkrete Zusagen, die im Gespräch gemacht wurden (von wem, was), als JSON-Array (leeres Array wenn keine)
4. naechste_schritte: offene Punkte / vereinbarte nächste Schritte als JSON-Array (leeres Array wenn keine)

Antworte NUR als JSON, keine Erklärung. Format:
{{"teilnehmer": [...], "kernpunkte": [...], "zusagen": [...], "naechste_schritte": [...]}}"""

    try:
        resp = client.messages.create(
            model="claude-sonnet-5",
            max_tokens=800,
            messages=[{"role": "user", "content": prompt}]
        )
        raw = resp.content[0].text.strip()
        raw = raw.replace("```json", "").replace("```", "").strip()
        return json.loads(raw)
    except Exception:
        return None


def process_file(filepath):
    if filepath.suffix.lower() in SKIP_EXTENSIONS:
        return False, "Code-Datei übersprungen"

    content = extract_text(filepath)
    if content is None:
        return False, "Extraktion fehlgeschlagen"

    result = classify(filepath, content)
    if not result:
        return False, "API-Klassifizierung fehlgeschlagen"

    zielordner = VAULT_PATH / result.get("zielordner", "Memos")
    zielordner.mkdir(parents=True, exist_ok=True)

    datum = datetime.now().strftime("%Y-%m-%d")

    # Nur für Text-Dokumente eine MD-Notiz anlegen (nicht für Medien)
    if filepath.suffix.lower() not in [".png", ".jpg", ".jpeg", ".gif", ".mp4", ".mov", ".mp3", ".zip"]:
        notiz_name = f"{datum}-{filepath.stem}.md"
        notiz_path = zielordner / notiz_name
        tags = result.get("tags", [])
        tags_str = "\n".join(f"  - {t}" for t in tags)
        volltext = content[:6000] if content else ""

        meeting_sections = ""
        if "/Meetings" in result.get("zielordner", ""):
            meeting_data = extract_meeting_structure(filepath)
            if meeting_data:
                def _liste(items):
                    return "\n".join(f"- {i}" for i in items) if items else "- (keine)"
                meeting_sections = f"""
## Teilnehmer
{_liste(meeting_data.get("teilnehmer", []))}

## Kernpunkte
{_liste(meeting_data.get("kernpunkte", []))}

## Zusagen
{_liste(meeting_data.get("zusagen", []))}

## Nächste Schritte
{_liste(meeting_data.get("naechste_schritte", []))}
"""

        notiz_inhalt = f"""---
tags:
{tags_str}
quelle: {filepath.name}
datum: {datum}
kategorie: {result.get("kategorie", "")}
---

# {filepath.stem}

## Zusammenfassung
{result.get("zusammenfassung", "")}
{meeting_sections}
## Vollständiger Inhalt
{volltext}
"""
        notiz_path.write_text(notiz_inhalt, encoding="utf-8")

    safe_name = filepath.name.replace("@", "-at-")
    ziel_original = zielordner / safe_name
    ziel = ziel_original if not ziel_original.exists() else zielordner / f"{datum}-{safe_name}"
    try:
        shutil.move(str(filepath), str(ziel))
    except Exception:
        try:
            shutil.copy2(str(filepath), str(ziel))
            filepath.unlink(missing_ok=True)
        except Exception:
            pass

    log_path = LOGS_PATH / "inbox_log.md"
    with open(log_path, "a", encoding="utf-8") as f:
        f.write(f"- {datetime.now().strftime('%Y-%m-%d %H:%M')} | {filepath.name} → {result.get('zielordner')} | {result.get('zusammenfassung', '')[:80]}\n")

    return True, result.get("zielordner", "")

def run():
    if not API_KEY:
        print("FEHLER: ANTHROPIC_API_KEY nicht gesetzt")
        return

    SKIP_NAMES = {"Thumbs.db", ".DS_Store", "package-lock.json", "package.json",
                  "tsconfig.json", "tsdoc-metadata.json", "openai"}
    SKIP_PREFIXES = ("README", "readme", "LICENSE", "license", "CHANGELOG",
                     "HISTORY", "History", "CONTRIBUTING", "contributing",
                     "~$")  # Word temp-Dateien (z.B. ~$sprechung.docx)

    dateien = [f for f in INBOX_PATH.rglob("*")
               if f.is_file()
               and not f.name.startswith(".")
               and not f.name.startswith("~$")
               and f.suffix.lower() not in SKIP_EXTENSIONS
               and f.name not in SKIP_NAMES
               and not any(f.name.startswith(p) for p in SKIP_PREFIXES)
               and "node_modules" not in str(f)
               and "_fehler" not in str(f)]

    cache = load_cache()
    dateien = [d for d in dateien if str(d) not in cache]

    if not dateien:
        print("Inbox leer (alle Dateien bereits verarbeitet).")
        return

    print(f"{len(dateien)} neue Datei(en) gefunden.")
    verarbeitet = 0
    fehler = 0

    for i, datei in enumerate(dateien, 1):
        print(f"  [{i}/{len(dateien)}] {datei.name} ...", end=" ", flush=True)
        try:
            ok, info = process_file(datei)
        except Exception as e:
            ok, info = False, f"Unerwarteter Fehler: {e}"
        if ok:
            cache.add(str(datei))
            save_cache(cache)
            print(f"→ {info}")
            verarbeitet += 1
        else:
            print(f"FEHLER: {info}")
            if info != "Code-Datei übersprungen":
                fehler_ziel = FEHLER_PATH / datei.name
                try:
                    shutil.move(str(datei), str(fehler_ziel))
                except Exception:
                    pass
            fehler += 1

    daily_file = DAILY_PATH / f"{datetime.now().strftime('%Y-%m-%d')}.md"
    with open(daily_file, "a", encoding="utf-8") as f:
        f.write(f"\n## Inbox-Verarbeitung {datetime.now().strftime('%H:%M')}\n")
        f.write(f"- Verarbeitet: {verarbeitet}\n")
        f.write(f"- Fehler: {fehler}\n")

    print(f"\nFertig: {verarbeitet} verarbeitet, {fehler} Fehler.")

if __name__ == "__main__":
    run()
